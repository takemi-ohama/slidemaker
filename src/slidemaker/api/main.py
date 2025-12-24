"""FastAPI main application for Slidemaker API.

このモジュールはSlidemakerのRESTful APIエンドポイントを提供します。
非同期タスク実行、ワークフロー統合、包括的なエラーハンドリングを実装しています。

Endpoints:
    GET  /              - API情報
    GET  /health        - ヘルスチェック
    POST /api/create    - Markdownからスライド作成
    POST /api/convert   - PDF/画像からスライド変換
    GET  /api/tasks/{id} - タスクステータス取得

Architecture:
    FastAPI Application
         ↓
    Background Tasks (NewSlideWorkflow / ConversionWorkflow)
         ↓
    S3 Storage (TaskManager + PowerPoint output)

Environment Variables:
    ANTHROPIC_API_KEY: Anthropic API key (必須)
    S3_BUCKET_NAME: S3 bucket name (必須)
    AWS_REGION: AWS region (オプション、デフォルト: us-east-1)
    ALLOWED_ORIGINS: CORS allowed origins (オプション、デフォルト: *)
    LOG_LEVEL: ログレベル (オプション、デフォルト: INFO)

Example:
    # Development server
    $ uvicorn slidemaker.api.main:app --reload

    # Production (AWS Lambda)
    from mangum import Mangum
    handler = Mangum(app)
"""

import base64
import os
import traceback
from contextlib import asynccontextmanager
from datetime import UTC, datetime
from pathlib import Path
from typing import Any, Literal

from fastapi import BackgroundTasks, FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

from slidemaker.api.schemas.requests import ConvertSlideRequest, CreateSlideRequest
from slidemaker.api.schemas.responses import (
    HealthCheckResponse,
    TaskResponse,
    TaskStatusResponse,
)
from slidemaker.api.storage import S3Storage
from slidemaker.api.tasks import TaskManager
from slidemaker.core.models.slide_config import SlideConfig
from slidemaker.image_processing.analyzer import ImageAnalyzer
from slidemaker.image_processing.loader import ImageLoader
from slidemaker.image_processing.processor import ImageProcessor
from slidemaker.llm.manager import LLMManager
from slidemaker.llm.models import LLMConfig
from slidemaker.pptx.generator import PowerPointGenerator
from slidemaker.utils.file_manager import FileManager
from slidemaker.utils.logger import get_logger, setup_logger
from slidemaker.workflows.conversion import ConversionWorkflow
from slidemaker.workflows.new_slide import NewSlideWorkflow

# ログ設定
setup_logger(level=os.getenv("LOG_LEVEL", "INFO"), format="json")
logger = get_logger(__name__)

# グローバル変数（Lambda環境での再利用のため）
storage: S3Storage | None = None
task_manager: TaskManager | None = None
llm_manager: LLMManager | None = None
file_manager: FileManager | None = None


@asynccontextmanager
async def lifespan(app: FastAPI) -> Any:
    """FastAPIライフサイクル管理（起動・終了処理）.

    Args:
        app: FastAPIアプリケーションインスタンス

    Yields:
        None

    Raises:
        ValueError: 必須環境変数が設定されていない場合
    """
    global storage, task_manager, llm_manager, file_manager

    logger.info("Starting Slidemaker API", version="0.6.0")

    # 環境変数チェック
    bucket_name = os.getenv("S3_BUCKET_NAME")
    if not bucket_name:
        raise ValueError("S3_BUCKET_NAME environment variable is required")

    anthropic_key = os.getenv("ANTHROPIC_API_KEY")
    if not anthropic_key:
        raise ValueError("ANTHROPIC_API_KEY environment variable is required")

    # S3Storage初期化
    region = os.getenv("AWS_REGION", "us-east-1")
    storage = S3Storage(bucket_name=bucket_name, region=region)
    await storage.__aenter__()
    logger.info("S3Storage initialized", bucket=bucket_name, region=region)

    # TaskManager初期化
    task_manager = TaskManager(storage=storage)
    logger.info("TaskManager initialized")

    # LLMManager初期化
    composition_config = LLMConfig(
        type="api",
        provider="claude",
        model=os.getenv("ANTHROPIC_MODEL", "claude-3-5-sonnet-20241022"),
        api_key=anthropic_key,
    )
    llm_manager = LLMManager(composition_config=composition_config)
    logger.info("LLMManager initialized")

    # FileManager初期化（一時ディレクトリ使用）
    file_manager = FileManager(output_base_dir="/tmp/slidemaker_output")
    logger.info("FileManager initialized")

    logger.info("Slidemaker API startup completed")

    yield

    # シャットダウン処理
    logger.info("Shutting down Slidemaker API")
    if storage:
        await storage.__aexit__(None, None, None)
        logger.info("S3Storage closed")


# FastAPIアプリケーション初期化
app = FastAPI(
    title="Slidemaker API",
    description="AI-Powered PowerPoint Generator API",
    version="0.6.0",
    lifespan=lifespan,
)

# CORS設定
allowed_origins = os.getenv("ALLOWED_ORIGINS", "*").split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
logger.info("CORS configured", allowed_origins=allowed_origins)


# ========================================
# エンドポイント
# ========================================


@app.get("/")
async def root() -> dict[str, str]:
    """ルートエンドポイント.

    Returns:
        API情報
    """
    return {
        "message": "Slidemaker API v0.6.0",
        "docs": "/docs",
    }


@app.get("/health", response_model=HealthCheckResponse)
async def health_check() -> HealthCheckResponse:
    """ヘルスチェックエンドポイント.

    システムの健全性をチェックします。

    Returns:
        ヘルスチェック結果

    Note:
        - LLM接続: ANTHROPIC_API_KEYの存在確認
        - S3接続: S3_BUCKET_NAMEの存在確認とS3Storageの初期化状態
    """
    checks: dict[str, bool] = {}

    # LLM接続チェック
    checks["llm_available"] = os.getenv("ANTHROPIC_API_KEY") is not None

    # S3接続チェック
    checks["storage_available"] = (
        os.getenv("S3_BUCKET_NAME") is not None and storage is not None
    )

    # 全体ステータス
    status: Literal["ok", "degraded", "down"]
    if all(checks.values()):
        status = "ok"
    elif any(checks.values()):
        status = "degraded"
    else:
        status = "down"

    return HealthCheckResponse(
        status=status,
        version="0.6.0",
        timestamp=datetime.now(UTC),
        checks=checks,
    )


@app.post("/api/create", response_model=TaskResponse)
async def create_slide(
    request: CreateSlideRequest, background_tasks: BackgroundTasks
) -> TaskResponse:
    """Markdownからスライドを作成します（非同期）.

    Args:
        request: スライド作成リクエスト
        background_tasks: FastAPI BackgroundTasks

    Returns:
        タスクレスポンス（task_id含む）

    Raises:
        HTTPException: タスク作成に失敗した場合（500）
    """
    if task_manager is None:
        raise HTTPException(status_code=500, detail="TaskManager not initialized")

    try:
        # タスク作成
        task_id = await task_manager.create_task()
        logger.info("Task created for create endpoint", task_id=task_id)

        # バックグラウンドでワークフロー実行
        background_tasks.add_task(_run_create_workflow, task_id, request)

        # タスク情報を返す
        return TaskResponse(
            task_id=task_id,
            status="pending",
            message="Task created successfully. Processing will start shortly.",
            created_at=datetime.now(UTC),
            updated_at=datetime.now(UTC),
        )

    except Exception as e:
        logger.error("Failed to create task", error=str(e), exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to create task") from e


@app.post("/api/convert", response_model=TaskResponse)
async def convert_slide(
    request: ConvertSlideRequest, background_tasks: BackgroundTasks
) -> TaskResponse:
    """PDF/画像からスライドを変換します（非同期）.

    Args:
        request: スライド変換リクエスト
        background_tasks: FastAPI BackgroundTasks

    Returns:
        タスクレスポンス（task_id含む）

    Raises:
        HTTPException: タスク作成に失敗した場合（500）
    """
    if task_manager is None:
        raise HTTPException(status_code=500, detail="TaskManager not initialized")

    try:
        # タスク作成
        task_id = await task_manager.create_task()
        logger.info("Task created for convert endpoint", task_id=task_id)

        # バックグラウンドでワークフロー実行
        background_tasks.add_task(_run_conversion_workflow, task_id, request)

        # タスク情報を返す
        return TaskResponse(
            task_id=task_id,
            status="pending",
            message="Task created successfully. Processing will start shortly.",
            created_at=datetime.now(UTC),
            updated_at=datetime.now(UTC),
        )

    except Exception as e:
        logger.error("Failed to create task", error=str(e), exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to create task") from e


@app.get("/api/tasks/{task_id}", response_model=TaskStatusResponse)
async def get_task_status(task_id: str) -> TaskStatusResponse:
    """タスクステータスを取得します.

    Args:
        task_id: タスクID（UUID形式）

    Returns:
        タスクステータスレスポンス

    Raises:
        HTTPException:
            - 404: タスクが存在しない場合
            - 500: ステータス取得に失敗した場合
    """
    if task_manager is None:
        raise HTTPException(status_code=500, detail="TaskManager not initialized")

    try:
        status = await task_manager.get_task_status(task_id)
        return status

    except ValueError as e:
        # タスクが存在しない
        logger.warning("Task not found", task_id=task_id)
        raise HTTPException(status_code=404, detail=str(e)) from e

    except Exception as e:
        logger.error(
            "Failed to get task status", task_id=task_id, error=str(e), exc_info=True
        )
        raise HTTPException(status_code=500, detail="Failed to get task status") from e


# ========================================
# バックグラウンドタスク
# ========================================


async def _run_create_workflow(task_id: str, request: CreateSlideRequest) -> None:
    """NewSlideWorkflowを実行します（バックグラウンド）.

    Args:
        task_id: タスクID
        request: スライド作成リクエスト

    Note:
        エラーが発生した場合はTaskManager.set_task_error()を呼び出します。
    """
    if task_manager is None or storage is None or llm_manager is None or file_manager is None:
        logger.error("Required managers not initialized", task_id=task_id)
        return

    try:
        # ステータス更新: processing
        await task_manager.update_task_status(
            task_id=task_id,
            status="processing",
            message="Starting slide generation workflow",
            progress=0.0,
        )
        logger.info("Create workflow started", task_id=task_id)

        # 出力ファイル名を決定
        output_filename = request.output_filename or f"slides_{task_id[:8]}.pptx"
        if not output_filename.endswith(".pptx"):
            output_filename += ".pptx"

        # 一時ファイルパス
        temp_output_path = Path(f"/tmp/{output_filename}")

        # NewSlideWorkflowを実行
        workflow = NewSlideWorkflow(
            llm_manager=llm_manager,
            file_manager=file_manager,
        )

        # 一時Markdownファイルを作成
        temp_markdown_path = Path(f"/tmp/input_{task_id[:8]}.md")
        with open(temp_markdown_path, "w", encoding="utf-8") as f:
            f.write(request.content)

        result = await workflow.execute(
            input_data=temp_markdown_path,
            output_path=temp_output_path,
            generate_images=True,  # 常に画像生成有効
        )

        # 一時Markdownファイルを削除
        if temp_markdown_path.exists():
            temp_markdown_path.unlink()

        logger.info("Create workflow completed", task_id=task_id, output_path=str(result))

        # 進捗更新: 80%（S3アップロード前）
        await task_manager.update_task_status(
            task_id=task_id,
            status="processing",
            message="Uploading PowerPoint to S3",
            progress=0.8,
        )

        # 結果をS3にアップロード
        with open(result, "rb") as f:
            pptx_data = f.read()

        s3_key = f"outputs/{task_id}/{output_filename}"
        await storage.upload_file(
            file_data=pptx_data,
            key=s3_key,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        logger.info("PowerPoint uploaded to S3", task_id=task_id, s3_key=s3_key)

        # 署名付きURL生成（7日間有効）
        presigned_url = await storage.generate_presigned_url(s3_key, expiration=604800)

        # ファイルサイズとページ数を取得
        file_size = len(pptx_data)
        # python-pptxでページ数を取得
        from pptx import Presentation as PptxPresentation

        prs = PptxPresentation(str(result))
        page_count = len(prs.slides)

        # タスク完了
        await task_manager.set_task_result(
            task_id=task_id,
            output_url=presigned_url,
            output_filename=output_filename,
            file_size=file_size,
            page_count=page_count,
        )

        logger.info("Create workflow task completed successfully", task_id=task_id)

        # 一時ファイルを削除
        if temp_output_path.exists():
            temp_output_path.unlink()

    except Exception as e:
        logger.error(
            "Create workflow failed",
            task_id=task_id,
            error=str(e),
            traceback=traceback.format_exc(),
        )
        await task_manager.set_task_error(
            task_id=task_id,
            error_code="INTERNAL_ERROR",
            error_message=f"Workflow execution failed: {str(e)}",
            details={"traceback": traceback.format_exc()},
        )


async def _run_conversion_workflow(
    task_id: str, request: ConvertSlideRequest
) -> None:
    """ConversionWorkflowを実行します（バックグラウンド）.

    Args:
        task_id: タスクID
        request: スライド変換リクエスト

    Note:
        エラーが発生した場合はTaskManager.set_task_error()を呼び出します。
    """
    if task_manager is None or storage is None or llm_manager is None or file_manager is None:
        logger.error("Required managers not initialized", task_id=task_id)
        return

    temp_input_path: Path | None = None
    temp_output_path: Path | None = None

    try:
        # ステータス更新: processing
        await task_manager.update_task_status(
            task_id=task_id,
            status="processing",
            message="Starting conversion workflow",
            progress=0.0,
        )
        logger.info("Conversion workflow started", task_id=task_id)

        # Base64デコード
        try:
            file_data = base64.b64decode(request.file_data, validate=True)
        except Exception as e:
            raise ValueError(f"Invalid Base64 encoding: {e}") from e

        # 一時ファイルに保存
        file_extension = ".pdf" if request.file_type == "pdf" else ".png"
        temp_input_path = Path(f"/tmp/input_{task_id[:8]}{file_extension}")
        with open(temp_input_path, "wb") as f:
            f.write(file_data)

        logger.info(
            "Input file saved",
            task_id=task_id,
            file_type=request.file_type,
            file_size=len(file_data),
        )

        # 設定を構築
        slide_config = request.config.to_slide_config() if request.config else None

        # 出力ファイル名を決定
        output_filename = request.output_filename or f"converted_{task_id[:8]}.pptx"
        if not output_filename.endswith(".pptx"):
            output_filename += ".pptx"

        # 一時出力ファイルパス
        temp_output_path = Path(f"/tmp/{output_filename}")

        # ConversionWorkflowを実行
        # デフォルトのSlideConfigを使用
        default_config = slide_config or SlideConfig()

        workflow = ConversionWorkflow(
            llm_manager=llm_manager,
            file_manager=file_manager,
            image_loader=ImageLoader(file_manager=file_manager),
            image_analyzer=ImageAnalyzer(llm_manager=llm_manager),
            image_processor=ImageProcessor(file_manager=file_manager),
            powerpoint_generator=PowerPointGenerator(config=default_config),
        )

        result = await workflow.execute(
            input_data=temp_input_path,
            output_path=temp_output_path,
        )

        logger.info("Conversion workflow completed", task_id=task_id, output_path=str(result))

        # 進捗更新: 80%（S3アップロード前）
        await task_manager.update_task_status(
            task_id=task_id,
            status="processing",
            message="Uploading PowerPoint to S3",
            progress=0.8,
        )

        # 結果をS3にアップロード
        with open(result, "rb") as f:
            pptx_data = f.read()

        s3_key = f"outputs/{task_id}/{output_filename}"
        await storage.upload_file(
            file_data=pptx_data,
            key=s3_key,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        logger.info("PowerPoint uploaded to S3", task_id=task_id, s3_key=s3_key)

        # 署名付きURL生成（7日間有効）
        presigned_url = await storage.generate_presigned_url(s3_key, expiration=604800)

        # ファイルサイズとページ数を取得
        file_size = len(pptx_data)
        # python-pptxでページ数を取得
        from pptx import Presentation as PptxPresentation

        prs = PptxPresentation(str(result))
        page_count = len(prs.slides)

        # タスク完了
        await task_manager.set_task_result(
            task_id=task_id,
            output_url=presigned_url,
            output_filename=output_filename,
            file_size=file_size,
            page_count=page_count,
        )

        logger.info("Conversion workflow task completed successfully", task_id=task_id)

        # 一時ファイルを削除
        if temp_input_path and temp_input_path.exists():
            temp_input_path.unlink()
        if temp_output_path and temp_output_path.exists():
            temp_output_path.unlink()

    except Exception as e:
        logger.error(
            "Conversion workflow failed",
            task_id=task_id,
            error=str(e),
            traceback=traceback.format_exc(),
        )
        await task_manager.set_task_error(
            task_id=task_id,
            error_code="INTERNAL_ERROR",
            error_message=f"Workflow execution failed: {str(e)}",
            details={"traceback": traceback.format_exc()},
        )

        # クリーンアップ
        if temp_input_path and temp_input_path.exists():
            temp_input_path.unlink()
        if temp_output_path and temp_output_path.exists():
            temp_output_path.unlink()


# ========================================
# エラーハンドラ
# ========================================


@app.exception_handler(ValueError)
async def value_error_handler(request: Any, exc: ValueError) -> JSONResponse:
    """ValueError用エラーハンドラ.

    Args:
        request: リクエストオブジェクト
        exc: 例外オブジェクト

    Returns:
        JSONレスポンス（400 Bad Request）
    """
    logger.warning("ValueError occurred", error=str(exc))
    return JSONResponse(status_code=400, content={"detail": str(exc)})


@app.exception_handler(FileNotFoundError)
async def not_found_handler(request: Any, exc: FileNotFoundError) -> JSONResponse:
    """FileNotFoundError用エラーハンドラ.

    Args:
        request: リクエストオブジェクト
        exc: 例外オブジェクト

    Returns:
        JSONレスポンス（404 Not Found）
    """
    logger.warning("FileNotFoundError occurred", error=str(exc))
    return JSONResponse(status_code=404, content={"detail": str(exc)})


@app.exception_handler(Exception)
async def general_exception_handler(request: Any, exc: Exception) -> JSONResponse:
    """汎用例外ハンドラ.

    Args:
        request: リクエストオブジェクト
        exc: 例外オブジェクト

    Returns:
        JSONレスポンス（500 Internal Server Error）
    """
    logger.error("Unhandled exception occurred", error=str(exc), exc_info=True)
    return JSONResponse(
        status_code=500,
        content={"detail": "Internal server error"},
    )
